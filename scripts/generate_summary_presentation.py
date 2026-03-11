from __future__ import annotations

import json
from collections import Counter
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = ROOT / "data"
EMBEDDINGS_DIR = DATA_DIR / "embeddings"
EXPORT_DIR = ROOT / "export"
ASSETS_DIR = EXPORT_DIR / "presentation_assets"
PRESENTATION_PATH = EXPORT_DIR / "ohbm2026-summary.pptx"

SLATE = "#1E293B"
BLUE = "#2563EB"
TEAL = "#0F766E"
ORANGE = "#EA580C"
GOLD = "#D97706"
GRAY = "#64748B"
LIGHT = "#E2E8F0"
RED = "#DC2626"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def ensure_dirs() -> None:
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)


def build_summary() -> dict:
    abstracts = load_json(DATA_DIR / "abstracts.json")["abstracts"]
    enriched = load_json(DATA_DIR / "abstracts_enriched.json")["abstracts"]
    authors = load_json(DATA_DIR / "authors.json")["authors"]
    references = load_json(DATA_DIR / "reference_metadata.json")
    image_analyses_path = DATA_DIR / "image_analyses.json"
    image_analyses = load_json(image_analyses_path) if image_analyses_path.exists() else {"analyses": []}

    accepted_for = Counter(abstract.get("accepted_for") for abstract in abstracts)
    primary_topics = Counter()
    for abstract in abstracts:
        for response in abstract.get("responses", []):
            if response.get("question_name") != "Primary Parent Category & Sub-Category":
                continue
            try:
                values = json.loads(response.get("value") or "[]")
            except json.JSONDecodeError:
                values = []
            if values:
                primary_topics[values[0]] += 1
            break

    matched_refs = sum(1 for item in references["references"] if item.get("openalex"))
    unmatched_refs = len(references["references"]) - matched_refs

    local_assets = [asset for abstract in abstracts for asset in abstract.get("local_assets", [])]
    downloaded_assets = sum(1 for asset in local_assets if asset.get("downloaded"))

    bundle_rows = []
    for meta_path in sorted(EMBEDDINGS_DIR.glob("*/metadata.json")):
        meta = load_json(meta_path)
        vectors = np.load(meta_path.parent / "vectors.npy", mmap_mode="r")
        bundle_rows.append(
            {
                "bundle": meta_path.parent.name,
                "model_name": meta.get("model_name") or meta.get("source_embedding_name") or "",
                "fields": meta.get("embedding_fields") or [],
                "shape": list(vectors.shape),
                "source": meta.get("source_embedding_name"),
            }
        )

    semantic_dirs = {
        "published_2": EMBEDDINGS_DIR / "voyage_stage2_published" / "semantic_analysis",
        "published_15": EMBEDDINGS_DIR / "voyage_stage2_published" / "semantic_analysis_15-communities",
        "published_21": EMBEDDINGS_DIR / "voyage_stage2_published" / "semantic_analysis_21-communities",
    }
    semantic = {}
    for key, path in semantic_dirs.items():
        summary = load_json(path / "community_detection.json")
        clusters = load_json(path / "cluster_summaries.json")["clusters"]
        semantic[key] = {"community": summary, "clusters": clusters}

    return {
        "abstract_count": len(abstracts),
        "enriched_count": len(enriched),
        "author_count": len(authors),
        "accepted_for": accepted_for,
        "primary_topics": primary_topics,
        "reference_unique_count": len(references["references"]),
        "reference_matched_count": matched_refs,
        "reference_unmatched_count": unmatched_refs,
        "reference_abstract_count": len(references["abstracts"]),
        "asset_total": len(local_assets),
        "asset_downloaded": downloaded_assets,
        "asset_missing": len(local_assets) - downloaded_assets,
        "asset_file_count": len(list((DATA_DIR / "assets").glob("*"))),
        "image_analysis_count": len(image_analyses.get("analyses", [])),
        "bundle_rows": bundle_rows,
        "semantic": semantic,
        "git_log": [
            "42ce05c Initial OHBM 2026 pipeline",
            "d8abefd Reorganize tests by module",
            "2e361b5 Document runtime requirements",
            "2b217d3 Break enrichment into task commands",
            "d056676 Simplify enriched abstract schema",
            "f377b43 Add local NeuroScape stage 2 pipeline",
            "eacc23a Add semantic visualization and OpenAlex tooling",
        ],
    }


def save_accepted_for_chart(summary: dict) -> Path:
    path = ASSETS_DIR / "accepted_for.png"
    labels = list(summary["accepted_for"].keys())
    values = [summary["accepted_for"][label] for label in labels]

    fig, ax = plt.subplots(figsize=(7, 4.2), dpi=180)
    bars = ax.bar(labels, values, color=[BLUE, ORANGE])
    ax.set_title("Accepted Presentation Types")
    ax.set_ylabel("Abstract count")
    ax.spines[["top", "right"]].set_visible(False)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 20, f"{value}", ha="center", va="bottom", fontsize=10)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_primary_topic_chart(summary: dict) -> Path:
    path = ASSETS_DIR / "primary_topics.png"
    items = summary["primary_topics"].most_common(8)
    labels = [item[0] for item in items][::-1]
    values = [item[1] for item in items][::-1]

    fig, ax = plt.subplots(figsize=(8, 5), dpi=180)
    ax.barh(labels, values, color=TEAL)
    ax.set_title("Top Primary Topics")
    ax.set_xlabel("Abstract count")
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_reference_chart(summary: dict) -> Path:
    path = ASSETS_DIR / "reference_coverage.png"
    values = [summary["reference_matched_count"], summary["reference_unmatched_count"]]
    labels = ["Matched", "Unmatched"]
    colors = [TEAL, LIGHT]

    fig, ax = plt.subplots(figsize=(5.5, 5.5), dpi=180)
    wedges, _ = ax.pie(values, colors=colors, startangle=90, wedgeprops={"width": 0.45, "edgecolor": "white"})
    ax.set_title("OpenAlex Reference Coverage")
    total = sum(values)
    for wedge, label, value in zip(wedges, labels, values):
        angle = (wedge.theta2 + wedge.theta1) / 2
        x = 0.8 * np.cos(np.deg2rad(angle))
        y = 0.8 * np.sin(np.deg2rad(angle))
        ax.text(x, y, f"{label}\n{value}\n{value/total:.0%}", ha="center", va="center", fontsize=10)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_embedding_family_chart(summary: dict) -> Path:
    path = ASSETS_DIR / "embedding_families.png"
    family_counts = Counter()
    for row in summary["bundle_rows"]:
        bundle = row["bundle"]
        if bundle.startswith("minilm"):
            family_counts["MiniLM"] += 1
        elif bundle.startswith("openai"):
            family_counts["OpenAI"] += 1
        elif bundle.startswith("pubmedbert"):
            family_counts["PubMedBERT"] += 1
        elif bundle.startswith("voyage"):
            family_counts["Voyage"] += 1
        elif bundle.startswith("neuroscape"):
            family_counts["NeuroScape stage 2"] += 1

    labels = list(family_counts.keys())
    values = [family_counts[label] for label in labels]
    colors = [BLUE, TEAL, ORANGE, GOLD, GRAY]

    fig, ax = plt.subplots(figsize=(7, 4.2), dpi=180)
    bars = ax.bar(labels, values, color=colors[: len(labels)])
    ax.set_title("Embedding Bundles by Family")
    ax.set_ylabel("Bundle count")
    ax.spines[["top", "right"]].set_visible(False)
    for bar, value in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, value + 0.05, str(value), ha="center", va="bottom", fontsize=10)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_umap_scatter() -> Path:
    path = ASSETS_DIR / "voyage_umap_primary_topic.png"
    data = load_json(EMBEDDINGS_DIR / "voyage_stage1" / "umap_title-introduction-methods-results-conclusion.json")
    points = data["points"]
    topic_counts = Counter(point["primary_topic"] for point in points)
    top_topics = [topic for topic, _ in topic_counts.most_common(8)]
    fallback = "Other"
    color_map = {topic: plt.cm.tab10(i % 10) for i, topic in enumerate(top_topics + [fallback])}

    fig, ax = plt.subplots(figsize=(8, 6), dpi=180)
    for topic in top_topics + [fallback]:
        topic_points = [
            point for point in points if (point["primary_topic"] if point["primary_topic"] in top_topics else fallback) == topic
        ]
        if not topic_points:
            continue
        ax.scatter(
            [point["x"] for point in topic_points],
            [point["y"] for point in topic_points],
            s=8,
            alpha=0.7,
            label=topic,
            color=color_map[topic],
        )
    ax.set_title("Voyage Stage 1 UMAP by Primary Topic")
    ax.set_xlabel("UMAP 1")
    ax.set_ylabel("UMAP 2")
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(loc="upper left", bbox_to_anchor=(1.02, 1.0), fontsize=8, frameon=False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def save_cluster_size_chart(summary: dict, key: str, filename: str, title: str) -> Path:
    path = ASSETS_DIR / filename
    clusters = summary["semantic"][key]["clusters"]
    items = sorted(((cluster["cluster_id"], cluster["size"], cluster["label"]) for cluster in clusters), key=lambda item: item[1], reverse=True)
    labels = [f"{cluster_id}: {label}" for cluster_id, _, label in items[:10]][::-1]
    values = [size for _, size, _ in items[:10]][::-1]

    fig, ax = plt.subplots(figsize=(8.5, 5.2), dpi=180)
    ax.barh(labels, values, color=BLUE if "15" in filename else ORANGE)
    ax.set_title(title)
    ax.set_xlabel("Abstract count")
    ax.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path


def set_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.7))
    paragraph = title_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(SLATE.replace("#", ""))
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.0), Inches(0.4))
        paragraph = subtitle_box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = subtitle
        run.font.size = Pt(13)
        run.font.color.rgb = RGBColor.from_string(GRAY.replace("#", ""))


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], font_size: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor.from_string(SLATE.replace("#", ""))


def add_image(slide, path: Path, left: float, top: float, width: float) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def add_table(slide, left: float, top: float, width: float, height: float, rows: list[list[str]]) -> None:
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for row_index, row in enumerate(rows):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(12 if row_index == 0 else 11)
                    run.font.bold = row_index == 0
                    run.font.color.rgb = RGBColor.from_string(SLATE.replace("#", ""))
            if row_index == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(LIGHT.replace("#", ""))


def create_presentation(summary: dict, chart_paths: dict[str, Path]) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "OHBM 2026 Abstract Pipeline Summary", "GraphQL ingest, enrichment, embeddings, OpenAlex references, and semantic clustering")
    add_bullets(
        slide,
        0.8,
        1.8,
        6.0,
        4.6,
        [
            f"{summary['abstract_count']:,} accepted abstracts ingested from Oxford Abstracts",
            f"{summary['author_count']:,} unique author records resolved",
            f"{summary['asset_downloaded']:,} figure files linked locally; {summary['asset_missing']} unresolved",
            f"{summary['reference_matched_count']:,}/{summary['reference_unique_count']:,} references matched to OpenAlex",
            "Multiple embedding families evaluated: MiniLM, PubMedBERT, OpenAI, Voyage, and NeuroScape stage 2",
        ],
    )
    add_bullets(
        slide,
        7.1,
        1.8,
        5.2,
        4.6,
        [
            "Repository activity summarized from seven main commits on `main`",
            "Interactive UMAP and UMAP vs t-SNE comparison exports generated for all embedding bundles",
            "Published NeuroScape stage 2 model applied to Voyage embeddings",
            "Community solutions produced at 2, 15, 21, and 27 clusters for semantic review",
        ],
        font_size=16,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "Ingest And Database Outputs")
    add_image(slide, chart_paths["accepted_for"], 0.7, 1.5, 5.4)
    add_image(slide, chart_paths["primary_topics"], 6.4, 1.4, 6.2)
    add_bullets(
        slide,
        0.8,
        5.8,
        12.0,
        1.0,
        [
            f"Local databases: abstracts.json ({summary['abstract_count']:,}), abstracts_enriched.json ({summary['enriched_count']:,}), authors.json ({summary['author_count']:,})"
        ],
        font_size=16,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "Figures And Reference Enrichment")
    add_image(slide, chart_paths["references"], 0.7, 1.5, 4.2)
    add_bullets(
        slide,
        5.2,
        1.7,
        7.0,
        3.8,
        [
            f"Methods/results figure policy enforced in the local asset table: {summary['asset_total']:,} figure entries, {summary['asset_downloaded']:,} downloaded, {summary['asset_missing']} unresolved",
            f"Cached figure understanding available for {summary['image_analysis_count']} figures using local Qwen multimodal analysis",
            f"OpenAlex reference metadata covers {summary['reference_matched_count']:,} unique references across {summary['reference_abstract_count']:,} abstracts",
            f"Unmatched references remaining after DOI/PMID/title search: {summary['reference_unmatched_count']:,}",
        ],
    )
    add_bullets(
        slide,
        5.2,
        5.6,
        7.0,
        0.8,
        ["Reference metadata output: data/reference_metadata.json"],
        font_size=15,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "Embedding Bundles Produced")
    add_image(slide, chart_paths["families"], 0.7, 1.5, 5.2)
    table_rows = [["Family", "Bundles", "Dimensions", "Field sets"]]
    table_rows.extend(
        [
            ["MiniLM", "4", "384", "default, methods/results, title/results/conclusion, title/introduction/conclusion"],
            ["PubMedBERT", "3", "768", "default, methods/results, title/results/conclusion"],
            ["OpenAI", "4", "1536", "default + 3 variants"],
            ["Voyage", "4", "1024", "default + 3 variants"],
            ["NeuroScape stage 2", "2", "64", "local MiniLM retrain and published Voyage transform"],
        ]
    )
    add_table(slide, 6.1, 1.5, 6.6, 4.2, table_rows)
    add_bullets(
        slide,
        6.2,
        5.95,
        6.2,
        0.7,
        ["All bundles include paired interactive UMAP and projection-comparison HTML/JSON outputs."],
        font_size=15,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "Voyage Stage 1 Geometry")
    add_image(slide, chart_paths["umap"], 0.8, 1.4, 11.8)
    add_bullets(
        slide,
        0.8,
        6.35,
        12.0,
        0.6,
        ["Representative export: voyage_stage1 UMAP colored by the dominant primary topics in the corpus."],
        font_size=15,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "Published NeuroScape Stage 2 On Voyage")
    add_bullets(
        slide,
        0.8,
        1.6,
        5.7,
        3.8,
        [
            "Published Zenodo stage 2 model applied to Voyage embeddings from the default field set",
            "Input/output: 1024-dimensional Voyage stage 1 -> 64-dimensional domain space",
            "Article similarity graph built with 3,333 nodes and 122,089 weighted edges",
            "Original unconstrained modularity preferred a coarse 2-cluster partition",
            "Fixed-resolution and minimum-community-count runs were added to inspect finer semantic structure",
        ],
    )
    add_table(
        slide,
        6.3,
        1.7,
        6.0,
        2.0,
        [
            ["Output", "Location"],
            ["Published stage 2 bundle", "data/embeddings/voyage_stage2_published"],
            ["2-cluster semantic output", "data/embeddings/voyage_stage2_published/semantic_analysis"],
            ["15-cluster semantic output", "data/embeddings/voyage_stage2_published/semantic_analysis_15-communities"],
            ["21-cluster semantic output", "data/embeddings/voyage_stage2_published/semantic_analysis_21-communities"],
        ],
    )
    add_bullets(
        slide,
        6.3,
        4.2,
        6.0,
        1.9,
        [
            "2-cluster modularity: 0.7539",
            "15-cluster modularity at resolution 2.5: 0.3958",
            "21-cluster modularity at resolution 2.9167: 0.3702",
        ],
        font_size=16,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "15-Cluster Semantic Solution")
    add_image(slide, chart_paths["clusters15"], 0.6, 1.4, 6.4)
    add_bullets(
        slide,
        7.3,
        1.5,
        5.3,
        4.8,
        [
            "Largest labels: brain/cortical/age; brain/pain/ad; connectivity/network/brain; task/language/visual; stimulation/network/motor",
            "This partition is the closest solution below the 16 primary-topic count and preserves stronger modularity than finer splits.",
            "Useful as a compact review layer for broad semantic neighborhoods before deeper drilling.",
        ],
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "21-Cluster Semantic Solution")
    add_image(slide, chart_paths["clusters21"], 0.6, 1.4, 6.4)
    add_bullets(
        slide,
        7.3,
        1.5,
        5.3,
        4.8,
        [
            "Largest labels: brain/using/data; pain/brain/ad; task/language/visual; connectivity/network/brain; brain/cortical/data",
            "This partition is the first solution above the 16 primary-topic count and exposes more specialized substructure.",
            "Useful as the starting point for manual topic review or downstream cluster annotation.",
        ],
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_title(slide, "Repository Activity And Deliverables")
    add_bullets(slide, 0.8, 1.5, 5.7, 4.8, summary["git_log"], font_size=16)
    add_bullets(
        slide,
        6.4,
        1.5,
        6.0,
        4.8,
        [
            "Key deliverables now on disk:",
            "abstracts.json, abstracts_enriched.json, authors.json, reference_metadata.json",
            "embedding bundles and interactive projection exports under data/embeddings and export/",
            "published-stage2 semantic outputs for 2, 15, 21, and 27 clusters",
            "CLI and tests updated to support reproducible fixed-resolution semantic analysis",
        ],
        font_size=16,
    )

    presentation.save(PRESENTATION_PATH)
    return PRESENTATION_PATH


def main() -> int:
    ensure_dirs()
    summary = build_summary()
    chart_paths = {
        "accepted_for": save_accepted_for_chart(summary),
        "primary_topics": save_primary_topic_chart(summary),
        "references": save_reference_chart(summary),
        "families": save_embedding_family_chart(summary),
        "umap": save_umap_scatter(),
        "clusters15": save_cluster_size_chart(
            summary,
            "published_15",
            "clusters_15.png",
            "Published Stage 2 Voyage: 15-Cluster Solution",
        ),
        "clusters21": save_cluster_size_chart(
            summary,
            "published_21",
            "clusters_21.png",
            "Published Stage 2 Voyage: 21-Cluster Solution",
        ),
    }
    path = create_presentation(summary, chart_paths)
    print(json.dumps({"presentation": str(path), "assets_dir": str(ASSETS_DIR)}, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
